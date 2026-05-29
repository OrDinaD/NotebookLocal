from __future__ import annotations

import logging
import zipfile
from pathlib import Path

from ingestion.chunking import chunk_text
from models.schemas import ImageAsset, IngestionError, IngestionResult, Modality, TableAsset, TextChunk

try:
    from docling.document_converter import DocumentConverter
except Exception:  # pragma: no cover - optional at test time
    DocumentConverter = None

try:
    import docx
except Exception:  # pragma: no cover
    docx = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    import openpyxl
except Exception:  # pragma: no cover
    openpyxl = None

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover
    PdfReader = None


logger = logging.getLogger(__name__)


class DocumentIngestor:
    def __init__(self, extract_dir: Path, chunk_size: int = 900, chunk_overlap: int = 150) -> None:
        self.extract_dir = extract_dir
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.extract_dir.mkdir(parents=True, exist_ok=True)

    def ingest(self, source_path: Path) -> IngestionResult:
        ext = source_path.suffix.lower()
        if ext in {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".txt", ".md"}:
            return self._ingest_document(source_path)
        if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"}:
            return self._ingest_image(source_path)
        return IngestionResult(
            source_uri=str(source_path),
            modality=Modality.document,
            errors=[IngestionError(code="unsupported_extension", message=f"Unsupported extension: {ext}")],
        )

    def _ingest_document(self, source_path: Path) -> IngestionResult:
        result = IngestionResult(source_uri=str(source_path), modality=Modality.document)
        docling_used = False

        if DocumentConverter is not None:
            try:
                converter = DocumentConverter()
                converted = converter.convert(str(source_path))
                text = converted.document.export_to_markdown()
                for idx, chunk in enumerate(chunk_text(text, self.chunk_size, self.chunk_overlap)):
                    result.text_chunks.append(
                        TextChunk(text=chunk, metadata={"source": str(source_path), "chunk_index": idx, "backend": "docling"})
                    )
                docling_used = True
            except Exception as exc:
                logger.warning("Docling failed for %s: %s", source_path, exc)
                result.errors.append(IngestionError(code="docling_failed", message=str(exc)))

        fallback = self._fallback_parse(source_path)
        if fallback.text_chunks:
            result.text_chunks = fallback.text_chunks if not result.text_chunks else result.text_chunks
        result.image_assets.extend(fallback.image_assets)
        result.table_assets.extend(fallback.table_assets)
        result.errors.extend(fallback.errors)

        if not docling_used and not result.text_chunks:
            result.errors.append(
                IngestionError(code="no_text_extracted", message="No text extracted via Docling or fallbacks")
            )

        return result

    def _ingest_image(self, source_path: Path) -> IngestionResult:
        text = f"Image file: {source_path.name}"
        return IngestionResult(
            source_uri=str(source_path),
            modality=Modality.image,
            text_chunks=[TextChunk(text=text, metadata={"source": str(source_path), "kind": "image_stub"})],
            image_assets=[ImageAsset(uri=str(source_path), metadata={"source": str(source_path)})],
        )

    def _fallback_parse(self, source_path: Path) -> IngestionResult:
        ext = source_path.suffix.lower()
        result = IngestionResult(source_uri=str(source_path), modality=Modality.document)

        if ext in {".txt", ".md"}:
            text = source_path.read_text(encoding="utf-8", errors="ignore")
            for idx, chunk in enumerate(chunk_text(text, self.chunk_size, self.chunk_overlap)):
                result.text_chunks.append(TextChunk(text=chunk, metadata={"chunk_index": idx, "backend": "plain_text"}))
            return result

        if ext == ".pdf":
            if PdfReader is None:
                result.errors.append(IngestionError(code="missing_dependency", message="pypdf is not installed"))
                return result
            reader = PdfReader(str(source_path))
            for page_idx, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                for chunk_idx, chunk in enumerate(chunk_text(page_text, self.chunk_size, self.chunk_overlap)):
                    result.text_chunks.append(
                        TextChunk(
                            text=chunk,
                            metadata={"page": page_idx, "chunk_index": chunk_idx, "backend": "pypdf"},
                        )
                    )
            return result

        if ext == ".docx":
            if docx is None:
                result.errors.append(IngestionError(code="missing_dependency", message="python-docx is not installed"))
                return result
            doc = docx.Document(str(source_path))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            for idx, chunk in enumerate(chunk_text(text, self.chunk_size, self.chunk_overlap)):
                result.text_chunks.append(TextChunk(text=chunk, metadata={"chunk_index": idx, "backend": "python-docx"}))
            result.image_assets.extend(self._extract_docx_images_zip(source_path))
            return result

        if ext == ".pptx":
            if Presentation is None:
                result.errors.append(IngestionError(code="missing_dependency", message="python-pptx is not installed"))
                return result
            pres = Presentation(str(source_path))
            for slide_idx, slide in enumerate(pres.slides, start=1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        txt = shape.text.strip()
                        if txt:
                            slide_texts.append(txt)
                    if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_bytes = shape.image.blob
                        out_path = self._dump_blob(source_path, slide_idx, len(result.image_assets), image_bytes, shape.image.ext)
                        result.image_assets.append(ImageAsset(uri=str(out_path), metadata={"slide": slide_idx}))
                merged = "\n".join(slide_texts)
                for chunk_idx, chunk in enumerate(chunk_text(merged, self.chunk_size, self.chunk_overlap)):
                    result.text_chunks.append(
                        TextChunk(
                            text=chunk,
                            metadata={"slide": slide_idx, "chunk_index": chunk_idx, "backend": "python-pptx"},
                        )
                    )
            return result

        if ext in {".xlsx", ".xlsm"}:
            if openpyxl is None:
                result.errors.append(IngestionError(code="missing_dependency", message="openpyxl is not installed"))
                return result
            wb = openpyxl.load_workbook(str(source_path), data_only=True)
            for ws in wb.worksheets:
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                    if values:
                        rows.append(" | ".join(values))
                sheet_text = "\n".join(rows)
                for chunk_idx, chunk in enumerate(chunk_text(sheet_text, self.chunk_size, self.chunk_overlap)):
                    result.text_chunks.append(
                        TextChunk(
                            text=chunk,
                            metadata={"sheet": ws.title, "chunk_index": chunk_idx, "backend": "openpyxl"},
                        )
                    )
                if rows:
                    result.table_assets.append(TableAsset(markdown="\n".join(rows[:20]), metadata={"sheet": ws.title}))
            return result

        result.errors.append(IngestionError(code="unsupported_fallback", message=f"No fallback parser for {ext}"))
        return result

    def _extract_docx_images_zip(self, source_path: Path) -> list[ImageAsset]:
        out: list[ImageAsset] = []
        target_dir = self.extract_dir / source_path.stem / "docx_media"
        target_dir.mkdir(parents=True, exist_ok=True)
        try:
            with zipfile.ZipFile(source_path, "r") as zip_ref:
                names = [n for n in zip_ref.namelist() if n.startswith("word/media/")]
                for idx, name in enumerate(names):
                    payload = zip_ref.read(name)
                    suffix = Path(name).suffix or ".bin"
                    out_path = target_dir / f"img_{idx}{suffix}"
                    out_path.write_bytes(payload)
                    out.append(ImageAsset(uri=str(out_path), metadata={"backend": "docx_zip", "index": idx}))
        except Exception as exc:
            logger.warning("Failed extracting DOCX images from %s: %s", source_path, exc)
        return out

    def _dump_blob(self, source_path: Path, slide_idx: int, image_idx: int, blob: bytes, ext: str) -> Path:
        target_dir = self.extract_dir / source_path.stem / "pptx_media"
        target_dir.mkdir(parents=True, exist_ok=True)
        out_path = target_dir / f"slide_{slide_idx:03d}_img_{image_idx:03d}.{ext or 'bin'}"
        out_path.write_bytes(blob)
        return out_path
